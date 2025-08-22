# _*_coding:utf-8 _*_

import traceback
import logging


'''文件解析模块(格式txt py .c, .cpp, .h .sh 等可以用txt打开的所有文件)'''


'''日志'''

logger = logging.getLogger(__name__)


'''文件打开函数'''

def read_file(filedir, encoding='utf-8', line=''):  # encoding 打开文件时所用的编码格式，默认utf-8
    try:
        # 获取文件line行标，有值则是按行读取，无则一次读取，一般100M以内直接读取，以上按行读取
        if line:
            logger.warning(f'按行读取文件，行数={line}')  # '开放中'
            return ''
        else:
            logger.warning(f'读取文件，文件路径={filedir}')
            with open(filedir, 'r', encoding=encoding) as data:
                return data.read()
    except Exception as ek:
        logger.error("打开文件错误:")
        logger.error(ek)
        logger.error(traceback.format_exc())
        return ''


'''pdf文件解析打开函数'''


import pdfplumber

# 常规打开方式（适合小文件，大于100M的使用分页方式读取）

def read_pdf(filedir, page2=''):
    try:
        # page2，有值则是按页读取，无则一次读取，一般100M以内直接读取，以上按页读取
        if page2:
            logger.warning(f'按页读取文件，行数={page2}')  # '开放中'
            return ''
        else:
            logger.warning(f'读取文件，文件路径={filedir}')
            pdflist = []
            with pdfplumber.open(filedir) as pdf:
                for page in pdf.pages:
                    pdflist.append(page.extract_text())

            # 返回数据，每个页面为一个元素
            return pdflist
    except Exception as ek:
        logger.error("打开文件错误:")
        logger.error(ek)
        logger.error(traceback.format_exc())
        return ''




'''docx文件打开python-docx'''


from docx import Document


def read_docx(file_path, page2=''):
    try:
        # page2，有值则是按页读取，无则一次读取，一般100M以内直接读取，以上按页读取
        if page2:
            logger.warning(f'按页读取文件，行数={page2}')  # '开放中'
            return ''
        else:
            logger.warning(f'读取文件，文件路径={file_path}')
            with open(file_path, 'rb') as file:
                doc = Document(file)
                full_text = []
                for paragraph in doc.paragraphs:
                    full_text.append(paragraph.text)

                # 返回文本列表
                return full_text
    except Exception as e:
        logger.error("打开docx文件错误:")
        logger.error(e)
        logger.error(traceback.format_exc())
        return ''



'''读取xlsx/xlsm/xltx/xltm文件'''


from openpyxl import load_workbook


def read_excel(file_path):
    try:
        # 加载工作簿
        workbook = load_workbook(filename=file_path)
        try:
            # 获取所有工作表
            sheets = workbook.sheetnames
            logger.warning(f"文件包含的工作表: {sheets}")

            data = {}
            # 读取每个工作表的内容
            for sheet_name in sheets:
                sheet = workbook[sheet_name]
                logger.warning(f"\n工作表 '{sheet_name}' 内容:")

                # 把sheet加到数据中
                if sheet_name not in data:
                    data[sheet_name] = []
                # 读取所有数据
                for row in sheet.iter_rows(values_only=True):
                    data[sheet_name].append(row)

            # 返回数据
            return data
        except Exception as e2:
            logger.error(f"读取工作表 '{file_path}' 时出错: {e2}")
            logger.error(traceback.format_exc())
            return ''
        finally:
            workbook.close()
    except Exception as e:
        logger.error(f"读取文件 '{file_path}' 时出错: {e}")
        logger.error(traceback.format_exc())
        return ''



'''读取ppt  python-pptx'''

from pptx import Presentation

def read_ppt(ppt_file_path):
    """读取PPT文件内容并返回文本"""
    try:
        logger.warning(f'ppt读取文件，文件路径={ppt_file_path}')
        with open(ppt_file_path, 'rb') as f:
            prs = Presentation(f)

            content = []
            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                content.append("\n".join(slide_content))

            logger.warning(f"PPT文件读取成功，现在返回内容")
            return content
    except Exception as e:
        logger.error(f"读取PPT文件时出错: {e}")
        logger.error(traceback.format_exc())
        return ''


'''读取csv文件'''

import csv
import chardet

def read_csv(file_path):
    try:
        logger.warning(f'csv读取文件，文件路径={file_path}')
        # 先检测文件编码
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
            # detected_encoding = result['encoding']
            # confidence = result['confidence']

        encoding = result['encoding'] if result['encoding'] else 'utf-8'
        print(f"检测到的编码: {encoding}")

        with open(file_path, 'r', encoding=encoding) as csvfile:
            # 创建csv阅读器
            csv_reader = csv.reader(csvfile)
            data = []
            for row in csv_reader:
                data.append(row)  # 每行是一个列表

            logger.warning(f"CSV文件读取成功，现在返回内容")
            return data
    except Exception as e:
        logger.error(f"读取CSV文件时出错: {e}")
        logger.error(traceback.format_exc())
        return []


'''读取docx并提取图片转为url'''

import os
import shutil
import zipfile
from docx.oxml.ns import qn
from xml.etree import ElementTree as ET
import time
import random

# 读取路径和url配置
with open('../file/conf.txt', 'r', encoding='utf-8') as data:
    conf = eval(data.read())
image_dir = conf.get('image_dir') if conf.get('image_dir') else "../file/public_img"
base_url = conf.get('base_url') if conf.get('base_url') else "https://example.com/images"

def read_docx_img(file_path, page2=''):
    try:
        """
        提取Word文档中的图片并替换为URL，返回修改后的文本内容

        参数:
        file_path: Word文档路径
        output_image_dir: 图片保存目录
        base_url: 图片基础URL (用于生成完整图片URL)

        返回:
        list: 替换图片为URL后的完整文本内容
        """
        # 拿到去掉后缀的文件名，创建保存图片的目录
        file_name = os.path.splitext(os.path.basename(file_path))[0]
        output_image_dir = f"{image_dir}/{file_name}"
        img_base_url = f"{base_url}{file_name}&filename="
        # 确保输出目录存在
        os.makedirs(output_image_dir, exist_ok=True)

        # 解压docx文件 (docx本质是zip)
        temp_dir = os.path.join(output_image_dir, "temp")
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)

        # 读取文档关系文件获取图片映射
        rels_path = os.path.join(temp_dir, "word", "_rels", "document.xml.rels")
        image_mapping = {}
        if os.path.exists(rels_path):
            tree = ET.parse(rels_path)
            root = tree.getroot()
            for child in root:
                if 'Relationship' in child.tag and 'image' in child.attrib['Type']:
                    r_id = child.attrib['Id']
                    image_path = os.path.join(temp_dir, "word", child.attrib['Target'])
                    image_mapping[r_id] = image_path

        # 处理文档内容
        doc = Document(file_path)
        all_text = []

        # 遍历文档所有段落
        for para in doc.paragraphs:
            para_text = ""
            runs_to_remove = []

            for run in para.runs:
                # 检查run中是否有图片
                if run._element.xpath('.//pic:pic', ):
                    # 获取图片关系ID
                    blip = run._element.xpath('.//a:blip')
                    if blip:
                        r_id = blip[0].get(qn('r:embed'))

                        if r_id in image_mapping:
                            # 保存图片
                            src_path = image_mapping[r_id]
                            ext = os.path.splitext(src_path)[1]
                            three_digit_random = random.randint(100, 999)
                            image_name = f"image_{int(time.time() * 1000)}{three_digit_random}{ext}"
                            dest_path = os.path.join(output_image_dir, image_name)
                            shutil.copy2(src_path, dest_path)

                            # 生成图片URL
                            if img_base_url:
                                image_url = f"{img_base_url}{image_name}"
                            else:
                                image_url = dest_path

                            # 在文本中标记图片位置
                            # position_tag = f"[{image_name.replace(ext, '')}]"
                            # para_text += position_tag

                            # 替换为URL
                            run.text = f' <img src="{image_url}" alt="{image_name.replace(ext, '')}"> '
                        else:
                            run.text = " [IMAGE_NOT_FOUND] "
                    else:
                        run.text = " [IMAGE_ERROR] "

                # 添加原始文本
                if run.text:
                    para_text += run.text

            all_text.append(para_text)

        # 清理临时文件
        shutil.rmtree(temp_dir)

        # 组合并返回完整文本
        return all_text
    except Exception as e:
        logger.error("打开docx文件错误:")
        logger.error(e)
        logger.error(traceback.format_exc())
        return ''














